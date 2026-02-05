"""Text and image extraction for various file formats (multimodal)."""

import base64
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Any

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".py",
    ".txt",
    ".md",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".csv",
    ".yaml",
    ".yml",
    ".ipynb",
}

# Image extensions supported for multimodal LLM input
IMAGE_EXTENSIONS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
    ".bmp",
}

# MIME types for image extensions
_IMAGE_MIME = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
}


def _bytes_to_data_url(raw: bytes, mime: str) -> str:
    """Encode raw image bytes as a data URL for LLM multimodal input."""
    b64 = base64.standard_b64encode(raw).decode("ascii")
    return f"data:{mime};base64,{b64}"


def _read_pdf(path: Path) -> str:
    """Extract text from PDF."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _read_docx(path: Path) -> str:
    """Extract text from Word document."""
    from docx import Document

    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _read_pptx(path: Path) -> str:
    """Extract text from PowerPoint."""
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if len(slide_texts) > 1:
            parts.append("\n".join(slide_texts))
    return "\n\n".join(parts)


def _read_text(path: Path) -> str:
    """Read plain text file with encoding fallback."""
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
    last_error = None
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError as e:
            last_error = e
    raise last_error or UnicodeDecodeError("unknown", b"", 0, 1, "")


def _read_ipynb(path: Path) -> str:
    """Convert Jupyter notebook to Python script using nbformat for grading.

    Extracts code cells and joins them into executable Python. Markdown and raw
    cells are excluded so the grader receives only the student's code.
    """
    import nbformat

    nb = nbformat.read(str(path), as_version=4)
    code_parts: list[str] = []
    for cell in nb.cells:
        if cell.cell_type != "code":
            continue
        source = cell.source
        if isinstance(source, list):
            source = "".join(source)
        if source.strip():
            code_parts.append(source.rstrip())
    return "\n\n".join(code_parts) if code_parts else ""


def _extract_images_from_ipynb(path: Path) -> list[tuple[bytes, str]]:
    """Extract embedded images from Jupyter notebook cell outputs.

    Collects image/png and image/jpeg from display_data and execute_result
    outputs so the grader can check plots and displayed images.
    """
    import nbformat

    result: list[tuple[bytes, str]] = []
    image_mimes = ("image/png", "image/jpeg", "image/gif", "image/webp")
    try:
        nb = nbformat.read(str(path), as_version=4)
        for cell in nb.cells:
            if cell.cell_type != "code" or not hasattr(cell, "outputs"):
                continue
            for out in cell.outputs:
                data = getattr(out, "data", None) or {}
                for mime in image_mimes:
                    b64_str = data.get(mime)
                    if not b64_str:
                        continue
                    if isinstance(b64_str, list):
                        b64_str = "".join(b64_str)
                    try:
                        raw = base64.standard_b64decode(b64_str)
                        result.append((raw, mime))
                    except Exception:
                        continue
    except Exception:
        pass
    return result


def _extract_images_from_pdf(path: Path) -> list[tuple[bytes, str]]:
    """Extract embedded images from a PDF. Returns list of (raw_bytes, mime_type)."""
    from pypdf import PdfReader

    result: list[tuple[bytes, str]] = []
    try:
        from PIL import Image  # noqa: F401 - needed for pypdf .image
    except ImportError:
        return result
    reader = PdfReader(path)
    for page in reader.pages:
        if not hasattr(page, "images"):
            continue
        try:
            names = list(page.images.keys()) if hasattr(page.images, "keys") else []
        except Exception:
            names = []
        for name in names:
            try:
                img_obj = page.images[name]
                pil = getattr(img_obj, "image", None)
                if pil is None:
                    continue
                buf = BytesIO()
                pil.save(buf, format="PNG")
                result.append((buf.getvalue(), "image/png"))
            except Exception:
                continue
    return result


def _extract_images_from_docx(path: Path) -> list[tuple[bytes, str]]:
    """Extract embedded images from a Word document. Returns list of (raw_bytes, mime_type)."""
    result: list[tuple[bytes, str]] = []
    try:
        with zipfile.ZipFile(path, "r") as zf:
            for name in zf.namelist():
                if not name.startswith("word/media/"):
                    continue
                ext = Path(name).suffix.lower()
                if ext not in _IMAGE_MIME:
                    continue
                try:
                    raw = zf.read(name)
                    result.append((raw, _IMAGE_MIME[ext]))
                except Exception:
                    continue
    except (zipfile.BadZipFile, OSError):
        pass
    return result


def _extract_images_from_pptx(path: Path) -> list[tuple[bytes, str]]:
    """Extract embedded images from a PowerPoint presentation.
    Returns list of (raw_bytes, mime_type)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    result: list[tuple[bytes, str]] = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                    continue
                img = getattr(shape, "image", None)
                if img is None:
                    continue
                blob = getattr(img, "blob", None)
                ext = (getattr(img, "ext", None) or "png").lower().strip()
                if not ext.startswith("."):
                    ext = "." + ext
                if blob and ext in _IMAGE_MIME:
                    result.append((blob, _IMAGE_MIME[ext]))
    except Exception:
        pass
    return result


def _read_image_as_data_url(path: Path) -> str:
    """Read image file and return a data URL (base64) for LLM multimodal input."""
    ext = path.suffix.lower()
    mime = _IMAGE_MIME.get(ext, "application/octet-stream")
    raw = path.read_bytes()
    b64 = base64.standard_b64encode(raw).decode("ascii")
    return f"data:{mime};base64,{b64}"


def extract_text_from_file(file_path: Path) -> str:
    """
    Extract text from a file based on its extension.

    Args:
        file_path: Path to the file.

    Returns:
        Extracted text content.

    Raises:
        ValueError: If the file format is not supported.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()

    if ext == ".pdf":
        return _read_pdf(path)
    if ext == ".docx":
        return _read_docx(path)
    if ext == ".pptx":
        return _read_pptx(path)
    if ext == ".ipynb":
        return _read_ipynb(path)
    if ext in SUPPORTED_EXTENSIONS:
        return _read_text(path)

    raise ValueError(f"Unsupported file format: {ext}")


def extract_content_parts_from_file(file_path: Path) -> list[dict[str, Any]]:
    """
    Extract content from a file as multimodal content parts for the LLM.

    - Text/code files: one "text" part.
    - Image files: file label + "image_url" part.
    - PDF/DOCX/PPTX: text + embedded images as "image_url" parts.

    Use with LangChain HumanMessage(content=...).

    Args:
        file_path: Path to the file.

    Returns:
        List of content parts (text and/or image_url).
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()
    file_label = f"=== FILE: {path.name} ===\n\n"

    if ext in IMAGE_EXTENSIONS:
        data_url = _read_image_as_data_url(path)
        return [
            {"type": "text", "text": file_label},
            {"type": "image_url", "image_url": {"url": data_url}},
        ]

    # Document types: text + embedded images
    if ext == ".pdf":
        text = _read_pdf(path)
        embedded = _extract_images_from_pdf(path)
    elif ext == ".docx":
        text = _read_docx(path)
        embedded = _extract_images_from_docx(path)
    elif ext == ".pptx":
        text = _read_pptx(path)
        embedded = _extract_images_from_pptx(path)
    elif ext == ".ipynb":
        text = _read_ipynb(path)
        embedded = _extract_images_from_ipynb(path)
    else:
        text = extract_text_from_file(path)
        embedded = []

    parts: list[dict[str, Any]] = [{"type": "text", "text": file_label + text}]
    for i, (raw, mime) in enumerate(embedded):
        label = f"\n[Image {i + 1} from {path.name}]\n"
        parts.append({"type": "text", "text": label})
        parts.append({"type": "image_url", "image_url": {"url": _bytes_to_data_url(raw, mime)}})
    return parts
